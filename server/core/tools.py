import json
import inspect
from typing import List, Dict, Callable

# Skill definitions

def create_word_doc(filename: str, content: str) -> str:
    """
    Create a Microsoft Word document (.docx).
    
    Args:
        filename (str): The name of the file to create (e.g., 'requirements.docx').
        content (str): The text content to write into the document.
        
    Returns:
        str: Success message or error.
    """
    try:
        from docx import Document
        import os
        
        # Ensure output directory exists
        out_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(out_dir, exist_ok=True)
        
        full_path = os.path.join(out_dir, filename)
        
        doc = Document()
        doc.add_heading(filename, 0)
        doc.add_paragraph(content)
        doc.save(full_path)
        return f"Document saved to {full_path}"
    except ImportError:
        return "Error: python-docx not installed."
    except Exception as e:
        return f"Error creating docx: {str(e)}"

def create_excel_sheet(filename: str, data_json: str) -> str:
    """
    Create a Microsoft Excel spreadsheet (.xlsx).
    
    Args:
        filename (str): The name of the file (e.g., 'data.xlsx').
        data_json (str): A JSON string representing a list of lists (rows).
                         Example: '[["Name", "Age"], ["Alice", 30], ["Bob", 25]]'
        
    Returns:
        str: Success message or error.
    """
    try:
        import openpyxl
        import os
        
        out_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(out_dir, exist_ok=True)
        full_path = os.path.join(out_dir, filename)
        
        wb = openpyxl.Workbook()
        ws = wb.active
        
        try:
            data = json.loads(data_json)
        except:
            return "Error: data_json must be a valid JSON string representing a list of lists."
            
        if isinstance(data, list):
            for row in data:
                if isinstance(row, list):
                    ws.append(row)
                else:
                    ws.append([str(row)])
        
        wb.save(full_path)
        return f"Spreadsheet saved to {full_path}"
    except ImportError:
        return "Error: openpyxl not installed."
    except Exception as e:
        return f"Error creating xlsx: {str(e)}"

def create_ppt(filename: str, slides_json: str) -> str:
    """
    Create a PowerPoint presentation (.pptx).
    
    Args:
        filename (str): The name of the file (e.g., 'slides.pptx').
        slides_json (str): A JSON string representing a list of slides.
                           Each slide is a dict with 'title' and 'content'.
                           Example: '[{"title": "Slide 1", "content": "Hello"}, ...]'
        
    Returns:
        str: Success message or error.
    """
    try:
        from pptx import Presentation
        import os
        
        out_dir = os.path.join(os.getcwd(), "output")
        os.makedirs(out_dir, exist_ok=True)
        full_path = os.path.join(out_dir, filename)
        
        prs = Presentation()
        
        try:
            slides_data = json.loads(slides_json)
        except:
            return "Error: slides_json must be a valid JSON string."
            
        if isinstance(slides_data, list):
            for slide_info in slides_data:
                slide_layout = prs.slide_layouts[1] # Bullet slide
                slide = prs.slides.add_slide(slide_layout)
                shapes = slide.shapes
                
                title_shape = shapes.title
                body_shape = shapes.placeholders[1]
                
                title_shape.text = slide_info.get("title", "Untitled")
                body_shape.text = slide_info.get("content", "")
                
        prs.save(full_path)
        return f"Presentation saved to {full_path}"
    except ImportError:
        return "Error: python-pptx not installed."
    except Exception as e:
        return f"Error creating pptx: {str(e)}"

# Registry
AVAILABLE_TOOLS = {
    "create_word_doc": create_word_doc,
    "create_excel_sheet": create_excel_sheet,
    "create_ppt": create_ppt
}

def get_tool_descriptions() -> str:
    """Generate a string describing available tools for the system prompt."""
    desc = "Available Tools:\n"
    for name, func in AVAILABLE_TOOLS.items():
        doc = inspect.getdoc(func)
        desc += f"- {name}: {doc}\n"
    return desc

def execute_tool(tool_name: str, **kwargs) -> str:
    """Execute a tool by name with arguments."""
    if tool_name in AVAILABLE_TOOLS:
        try:
            return AVAILABLE_TOOLS[tool_name](**kwargs)
        except Exception as e:
            return f"Error executing {tool_name}: {str(e)}"
    return f"Tool {tool_name} not found."
